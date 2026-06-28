"""
AI-Powered Playwright Framework — Professional PowerPoint Generator
Run: python3 docs/generate-presentation.py
Output: docs/AI-Platform-Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Brand colours ─────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x1b, 0x2a)   # deep navy
BG_CARD  = RGBColor(0x12, 0x24, 0x36)   # slightly lighter navy
ACCENT1  = RGBColor(0x58, 0xa6, 0xff)   # blue
ACCENT2  = RGBColor(0x3f, 0xb9, 0x50)   # green
ACCENT3  = RGBColor(0xf0, 0x88, 0x3e)   # orange
ACCENT4  = RGBColor(0xbc, 0x8c, 0xff)   # purple
ACCENT5  = RGBColor(0xf8, 0x51, 0x49)   # red
WHITE    = RGBColor(0xe6, 0xed, 0xf3)   # near-white
MUTED    = RGBColor(0x8b, 0x94, 0x9e)   # muted grey
BORDER   = RGBColor(0x30, 0x36, 0x3d)   # border grey
DARK_BG  = RGBColor(0x16, 0x1b, 0x22)   # card bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_pill(slide, text, left, top, fill_color, text_color):
    w = Inches(1.6)
    h = Inches(0.3)
    rect = add_rect(slide, left, top, w, h, fill_color, None)
    tf = rect.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color
    return rect

def add_divider(slide, top):
    line = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

def eyebrow(slide, text):
    add_textbox(slide, text,
                Inches(0.6), Inches(0.22), Inches(11), Inches(0.4),
                font_size=9, bold=True, color=ACCENT1)

def slide_title(slide, text, top=Inches(0.65)):
    add_textbox(slide, text,
                Inches(0.6), top, Inches(12.1), Inches(0.8),
                font_size=28, bold=True, color=WHITE)

def body_text(slide, text, left, top, width, height, size=13, color=MUTED):
    add_textbox(slide, text, left, top, width, height,
                font_size=size, color=color, wrap=True)

def bullet_item(slide, number, text, top, color=ACCENT1):
    # Number badge
    badge = add_rect(slide, Inches(0.6), top + Inches(0.04),
                     Inches(0.28), Inches(0.28), color, None)
    btf = badge.text_frame
    bp  = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br  = bp.add_run()
    br.text = str(number)
    br.font.size  = Pt(10)
    br.font.bold  = True
    br.font.color.rgb = WHITE
    # Text
    add_textbox(slide, text,
                Inches(1.0), top, Inches(5.5), Inches(0.36),
                font_size=12, color=WHITE, wrap=True)

def card(slide, left, top, width, height, border_color=BORDER, fill=DARK_BG):
    return add_rect(slide, left, top, width, height, fill, border_color, Pt(1.2))

def card_with_text(slide, left, top, width, height,
                   label, label_color, body_lines,
                   border_color=BORDER):
    card(slide, left, top, width, height, border_color)
    add_textbox(slide, label,
                left + Inches(0.18), top + Inches(0.12),
                width - Inches(0.36), Inches(0.28),
                font_size=9, bold=True, color=label_color)
    add_textbox(slide, body_lines,
                left + Inches(0.18), top + Inches(0.44),
                width - Inches(0.36), height - Inches(0.55),
                font_size=11, color=MUTED, wrap=True)

# ── Blank layout (index 6 is usually blank) ───────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, BG)
    return slide

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
def slide_01_title(prs):
    s = blank_slide(prs)
    # Large accent bar on left
    add_rect(s, Inches(0), Inches(1.5), Inches(0.08), Inches(4.5), ACCENT1)

    add_textbox(s, "AI TEST INTELLIGENCE PLATFORM",
                Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
                font_size=11, bold=True, color=ACCENT1)

    add_textbox(s, "AI-Powered Playwright\nTest Generation Framework",
                Inches(0.6), Inches(2.1), Inches(11), Inches(1.6),
                font_size=40, bold=True, color=WHITE, wrap=True)

    add_textbox(s,
        "Plain-English requirements → Production-quality Playwright tests.\n"
        "Automatic. Typed. CI-ready. Self-healing.",
                Inches(0.6), Inches(3.85), Inches(10), Inches(0.8),
                font_size=16, color=MUTED, wrap=True)

    # Pills row
    pills = [
        ("TypeScript Strict", ACCENT1, RGBColor(0x1f,0x3c,0x6e)),
        ("Playwright + Allure", ACCENT2, RGBColor(0x0f,0x2d,0x1a)),
        ("Multi-Provider LLM", ACCENT3, RGBColor(0x3d,0x20,0x00)),
        ("9 AI Modules", ACCENT4, RGBColor(0x2d,0x1f,0x6e)),
        ("Self-Healing Tests", ACCENT2, RGBColor(0x0f,0x2d,0x1a)),
    ]
    x = Inches(0.6)
    for label, tc, fc in pills:
        p = add_pill(s, label, x, Inches(5.0), fc, tc)
        x += Inches(1.7)

    add_textbox(s, "Ravi Gaygol  ·  2026",
                Inches(0.6), Inches(6.6), Inches(6), Inches(0.4),
                font_size=11, color=MUTED)
    add_textbox(s, "automationexercise.com",
                Inches(9.5), Inches(6.6), Inches(3.3), Inches(0.4),
                font_size=11, color=ACCENT1, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
def slide_02_exec_summary(prs):
    s = blank_slide(prs)
    eyebrow(s, "EXECUTIVE SUMMARY")
    slide_title(s, "What This Platform Does")
    add_divider(s, Inches(1.55))

    summary = (
        "This framework eliminates the manual work of writing Playwright test automation. "
        "A QA engineer provides a plain-English requirement; the AI pipeline produces a "
        "complete, typed, CI-ready spec file with page objects, test data, and assertions "
        "— in under two minutes."
    )
    add_textbox(s, summary,
                Inches(0.6), Inches(1.72), Inches(12.1), Inches(0.9),
                font_size=14, color=MUTED, wrap=True)

    stats = [
        ("< 2 min",    "Test generation time\nper requirement",      ACCENT1, RGBColor(0x1f,0x3c,0x6e)),
        ("104",        "Unit tests on the\nframework itself",        ACCENT2, RGBColor(0x0f,0x2d,0x1a)),
        ("0 errors",   "TypeScript strict mode\nacross all modules", ACCENT3, RGBColor(0x3d,0x20,0x00)),
        ("3 providers","Gemini · GitHub Models\n· OpenRouter",       ACCENT4, RGBColor(0x2d,0x1f,0x6e)),
    ]
    x = Inches(0.6)
    for val, label, tc, fc in stats:
        card(s, x, Inches(2.8), Inches(2.9), Inches(1.6), tc)
        add_textbox(s, val,
                    x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.65),
                    font_size=30, bold=True, color=tc)
        add_textbox(s, label,
                    x + Inches(0.2), Inches(3.6), Inches(2.5), Inches(0.6),
                    font_size=11, color=MUTED, wrap=True)
        x += Inches(3.05)

    points = [
        "Runs on automationexercise.com — a real, live web application",
        "Swappable LLM providers: no vendor lock-in",
        "Generated tests are indistinguishable from hand-written ones",
        "Self-healing, flaky analysis, root cause, coverage gap — all automated",
        "CI pipeline: TypeScript check → ESLint → Unit tests → Smoke tests",
    ]
    y = Inches(4.65)
    for i, pt in enumerate(points, 1):
        add_textbox(s, f"{'●'}  {pt}",
                    Inches(0.6), y, Inches(12), Inches(0.35),
                    font_size=12, color=WHITE)
        y += Inches(0.37)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
def slide_03_problem(prs):
    s = blank_slide(prs)
    eyebrow(s, "THE PROBLEM")
    slide_title(s, "Why Traditional Test Automation Breaks Down")
    add_divider(s, Inches(1.55))

    problems = [
        ("💸", "High Writing Cost",
         "QA engineers spend 60–70% of their time writing boilerplate — not designing strategies."),
        ("🔧", "Fragile Selectors",
         "One UI rename breaks dozens of tests. Maintenance cost outpaces new test creation."),
        ("🐌", "Slow CI Feedback",
         "Full regression runs on every PR even when 1 file changed. 45-minute pipelines block deploys."),
        ("🔍", "No Visibility",
         "No automated way to know which requirements lack coverage until a bug slips to production."),
        ("🩹", "Manual Debugging",
         "When tests fail flakily, engineers spend hours reading stack traces to find timing issues."),
        ("📋", "Copy-Paste Tests",
         "Test case ideas stay in spreadsheets. Converting them to code is a full day of work."),
    ]

    x_start = Inches(0.6)
    y_start = Inches(1.75)
    card_w   = Inches(3.9)
    card_h   = Inches(1.55)
    gap      = Inches(0.25)

    for i, (icon, title, desc) in enumerate(problems):
        col = i % 3
        row = i // 3
        x = x_start + col * (card_w + gap)
        y = y_start + row * (card_h + gap)
        card(s, x, y, card_w, card_h, BORDER)
        add_textbox(s, icon,
                    x + Inches(0.15), y + Inches(0.12),
                    Inches(0.5), Inches(0.4), font_size=18)
        add_textbox(s, title,
                    x + Inches(0.65), y + Inches(0.12),
                    card_w - Inches(0.8), Inches(0.38),
                    font_size=13, bold=True, color=WHITE)
        add_textbox(s, desc,
                    x + Inches(0.15), y + Inches(0.58),
                    card_w - Inches(0.3), Inches(0.85),
                    font_size=11, color=MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
def slide_04_solution(prs):
    s = blank_slide(prs)
    eyebrow(s, "THE SOLUTION")
    slide_title(s, "One Sentence In. Running Tests Out.")
    add_divider(s, Inches(1.55))

    # Input card
    card(s, Inches(0.6), Inches(1.72), Inches(5.5), Inches(1.0), ACCENT1)
    add_textbox(s, "INPUT",
                Inches(0.8), Inches(1.78), Inches(2), Inches(0.28),
                font_size=9, bold=True, color=ACCENT1)
    add_textbox(s, '"User should be able to login\nwith valid credentials"',
                Inches(0.8), Inches(2.05), Inches(5.1), Inches(0.55),
                font_size=13, color=WHITE, italic=True)

    # Arrow
    add_textbox(s, "⬇  AI Pipeline  ⬇",
                Inches(2.2), Inches(2.82), Inches(3), Inches(0.4),
                font_size=14, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)

    # Output card
    card(s, Inches(0.6), Inches(3.3), Inches(5.5), Inches(2.85), ACCENT2)
    add_textbox(s, "OUTPUT — tests/e2e/login.spec.ts",
                Inches(0.8), Inches(3.38), Inches(5.1), Inches(0.28),
                font_size=9, bold=True, color=ACCENT2)
    outputs = [
        "✓  8 test cases (positive + negative + edge)",
        "✓  Realistic test data (valid + invalid)",
        "✓  Real selectors from knowledge base",
        "✓  Correct assertions for every scenario",
        "✓  Ready to run with npm test",
    ]
    y = Inches(3.72)
    for line in outputs:
        add_textbox(s, line, Inches(0.8), y, Inches(5.1), Inches(0.34),
                    font_size=12, color=WHITE)
        y += Inches(0.36)

    # Right side — how it works
    steps = [
        (ACCENT1, "ExcelReader",        "Reads requirements from .xlsx"),
        (ACCENT2, "KnowledgeBase",      "Loads real selectors + page context"),
        (ACCENT3, "TestCaseGenerator",  "Requirement → 4–10 test cases"),
        (ACCENT4, "TestDataGenerator",  "Generates valid + invalid test data"),
        (ACCENT1, "ActionModelGen",     "Step → structured action JSON"),
        (ACCENT2, "PlaywrightRenderer", "Action JSON → TypeScript code"),
    ]
    y = Inches(1.72)
    for color, name, desc in steps:
        add_rect(s, Inches(6.5), y + Inches(0.06),
                 Inches(0.06), Inches(0.35), color)
        add_textbox(s, name,
                    Inches(6.7), y, Inches(3), Inches(0.32),
                    font_size=12, bold=True, color=WHITE)
        add_textbox(s, desc,
                    Inches(6.7), y + Inches(0.3), Inches(5.8), Inches(0.32),
                    font_size=10, color=MUTED)
        add_textbox(s, "↓",
                    Inches(6.7), y + Inches(0.62), Inches(1), Inches(0.25),
                    font_size=11, color=BORDER)
        y += Inches(0.86)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — FRAMEWORK ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide_05_architecture(prs):
    s = blank_slide(prs)
    eyebrow(s, "SYSTEM DESIGN")
    slide_title(s, "Four-Layer Architecture")
    add_divider(s, Inches(1.55))

    layers = [
        (ACCENT2, "1", "LLM Provider Layer",
         "The AI connection. Single LLMProvider interface → 4 implementations.",
         "GeminiProvider · GitHubModelsProvider · OpenRouterProvider · MockLLMProvider",
         "pipeline/providers/  ·  Swappable via LLM_PROVIDER env var"),
        (ACCENT3, "2", "Knowledge Base Layer",
         "Ground truth. Prevents hallucination. JSON files per page.",
         "KnowledgeBaseService · KnowledgeBaseGenerator · SelectorRetriever (RAG)",
         "pipeline/kb/pages/  ·  ae-home.json · ae-login.json"),
        (ACCENT1, "3", "AI Modules Layer",
         "Nine independent, independently testable intelligence modules.",
         "TestCaseGenerator · TestDataGenerator · ActionModelGen · AssertionGenerator · Self-Healing · Flaky · RootCause · Coverage · Regression",
         "pipeline/generators/  ·  pipeline/analyzers/"),
        (ACCENT4, "4", "Code Generation Layer",
         "No AI. Deterministic. Takes structured JSON, renders TypeScript.",
         "PlaywrightGenerator orchestrates  ·  PlaywrightRenderer converts ActionModel → code",
         "pipeline/generators/playwright/"),
    ]

    y = Inches(1.72)
    layer_h = Inches(1.22)
    for color, num, title, desc, modules, paths in layers:
        add_rect(s, Inches(0.6), y, Inches(0.06), layer_h, color)
        card(s, Inches(0.72), y, Inches(12.0), layer_h, color, DARK_BG)

        # Number
        r2 = int(color[0]) // 3
        g2 = int(color[1]) // 3
        b2 = int(color[2]) // 3
        add_rect(s, Inches(0.9), y + Inches(0.42),
                 Inches(0.36), Inches(0.36),
                 RGBColor(r2, g2, b2))
        add_textbox(s, num,
                    Inches(0.9), y + Inches(0.42), Inches(0.36), Inches(0.36),
                    font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

        add_textbox(s, title,
                    Inches(1.42), y + Inches(0.08), Inches(5), Inches(0.36),
                    font_size=14, bold=True, color=color)
        add_textbox(s, desc,
                    Inches(1.42), y + Inches(0.44), Inches(5.5), Inches(0.34),
                    font_size=11, color=WHITE)
        add_textbox(s, modules,
                    Inches(1.42), y + Inches(0.78), Inches(11.0), Inches(0.32),
                    font_size=10, color=MUTED, italic=True)
        add_textbox(s, paths,
                    Inches(7.5), y + Inches(0.08), Inches(5.0), Inches(0.36),
                    font_size=10, color=BORDER, align=PP_ALIGN.RIGHT)
        y += layer_h + Inches(0.06)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
def slide_06_workflow(prs):
    s = blank_slide(prs)
    eyebrow(s, "AI WORKFLOW")
    slide_title(s, "End-to-End: Requirement to Running Test")
    add_divider(s, Inches(1.55))

    steps = [
        (ACCENT1, "requirements.xlsx", "Fill Page, Feature, Description"),
        (ACCENT3, "ExcelReader", "Parses rows into Requirement objects"),
        (ACCENT2, "KB Loader", "Loads ae-home.json — real selectors"),
        (ACCENT1, "TestCaseGen", "4–10 test cases per requirement"),
        (ACCENT4, "TestDataGen", "Valid + invalid data sets"),
        (ACCENT3, "ActionModel", "Each step → structured JSON"),
        (ACCENT2, "Assertions", "Expected results → typed assertions"),
        (ACCENT1, "Playwright\nRenderer", "JSON → TypeScript code"),
        (ACCENT2, "tests/e2e/\n.spec.ts", "Ready to run"),
    ]

    box_w = Inches(1.35)
    box_h = Inches(1.0)
    y_box  = Inches(2.1)
    x = Inches(0.4)

    for i, (color, title, subtitle) in enumerate(steps):
        card(s, x, y_box, box_w, box_h, color, DARK_BG)
        add_textbox(s, title,
                    x + Inches(0.08), y_box + Inches(0.08),
                    box_w - Inches(0.16), Inches(0.52),
                    font_size=11, bold=True, color=color,
                    align=PP_ALIGN.CENTER, wrap=True)
        add_textbox(s, subtitle,
                    x + Inches(0.05), y_box + Inches(0.6),
                    box_w - Inches(0.1), Inches(0.36),
                    font_size=8, color=MUTED,
                    align=PP_ALIGN.CENTER, wrap=True)
        if i < len(steps) - 1:
            add_textbox(s, "→",
                        x + box_w + Inches(0.03), y_box + Inches(0.38),
                        Inches(0.25), Inches(0.3),
                        font_size=16, color=BORDER,
                        align=PP_ALIGN.CENTER)
        x += box_w + Inches(0.28)

    # Description below
    add_textbox(s,
        "Every AI module follows the same pattern:\n"
        "Input Interface  →  Engineered Prompt  →  LLM Call  →  AIJsonParser  →  Validated Typed Output",
                Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.65),
                font_size=12, color=MUTED, wrap=True)

    # Command
    card(s, Inches(0.6), Inches(4.35), Inches(6), Inches(0.55), BORDER, DARK_BG)
    add_textbox(s, "$ npm run ai:run",
                Inches(0.9), Inches(4.42), Inches(5.5), Inches(0.4),
                font_size=13, color=ACCENT2)

    card(s, Inches(6.8), Inches(4.35), Inches(6), Inches(0.55), BORDER, DARK_BG)
    add_textbox(s, "$ LLM_PROVIDER=gemini npm run ai:run",
                Inches(7.1), Inches(4.42), Inches(5.5), Inches(0.4),
                font_size=13, color=ACCENT2)

    # Provider info
    providers = [
        ("Gemini", "GOOGLE_API_KEY", ACCENT2),
        ("GitHub Models", "GITHUB_TOKEN", ACCENT1),
        ("OpenRouter", "OPENROUTER_API_KEY", ACCENT3),
        ("Mock (no cost)", "—", MUTED),
    ]
    x = Inches(0.6)
    for name, key, color in providers:
        add_textbox(s, f"◆  {name}",
                    x, Inches(5.15), Inches(2.9), Inches(0.3),
                    font_size=11, bold=True, color=color)
        add_textbox(s, key,
                    x, Inches(5.45), Inches(2.9), Inches(0.28),
                    font_size=10, color=MUTED)
        x += Inches(3.0)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FOLDER STRUCTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide_07_folder_structure(prs):
    s = blank_slide(prs)
    eyebrow(s, "PROJECT LAYOUT")
    slide_title(s, "Every Folder Has One Job")
    add_divider(s, Inches(1.55))

    folders = [
        (ACCENT1, "pipeline/",       "All AI logic lives here"),
        (ACCENT1, "  providers/",    "LLM interface + 4 implementations + fallback"),
        (ACCENT1, "  kb/",           "Knowledge base: JSON page files + RAG retriever"),
        (ACCENT1, "  generators/",   "9 generators: test-cases, action-model, pom, playwright…"),
        (ACCENT1, "  analyzers/",    "5 analyzers: flaky, root-cause, coverage, regression, healing"),
        (ACCENT3, "tests/",          "Standard Playwright layer"),
        (ACCENT3, "  e2e/",          "Spec files — generated + hand-written"),
        (ACCENT3, "  pages/",        "Page Object Models (AeHomePage, AeLoginPage)"),
        (ACCENT3, "  fixtures/",     "testDesktop + testMobile custom fixtures"),
        (ACCENT2, "scripts/",        "CLI entry points — npm run ai:run, demo:*, generate:*"),
        (ACCENT4, "config/",         "platform.json — suite definitions"),
        (ACCENT4, "requirements/",   "requirements.xlsx — input for Excel pipeline"),
    ]

    y = Inches(1.72)
    for color, name, desc in folders:
        is_parent = not name.startswith("  ")
        indent = Inches(0) if is_parent else Inches(0.3)
        if is_parent:
            add_rect(s, Inches(0.6) + indent, y + Inches(0.07),
                     Inches(0.04), Inches(0.26), color)
        add_textbox(s, name,
                    Inches(0.78) + indent, y,
                    Inches(3.0), Inches(0.36),
                    font_size=12 if is_parent else 11,
                    bold=is_parent,
                    color=color if is_parent else MUTED)
        add_textbox(s, desc,
                    Inches(4.2), y, Inches(8.8), Inches(0.36),
                    font_size=11, color=WHITE if is_parent else MUTED)
        y += Inches(0.36) if not is_parent else Inches(0.4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
def slide_08_features(prs):
    s = blank_slide(prs)
    eyebrow(s, "CAPABILITIES")
    slide_title(s, "What the Framework Does Out of the Box")
    add_divider(s, Inches(1.55))

    features = [
        ("⚡", ACCENT1, "AI Test Generation",
         "Requirement → 4–10 test cases covering positive, negative, edge, and security scenarios."),
        ("🗂️", ACCENT2, "Excel Input Pipeline",
         "QA teams fill requirements.xlsx; AI generates everything else automatically."),
        ("🔧", ACCENT3, "Auto Page Objects",
         "POMGenerator creates TypeScript POMs from knowledge base JSON. private readonly locators, public methods."),
        ("🩹", ACCENT5, "Self-Healing Locators",
         "Broken selector + page name → healed selector from KB with confidence score."),
        ("📊", ACCENT4, "Coverage Analysis",
         "Requirements vs existing tests → gap report. Know what is and isn't covered."),
        ("⚠️", ACCENT5, "Flaky Test Analyzer",
         "Execution metrics → flakiness probability + specific fix recommendation."),
        ("🔍", ACCENT3, "Root Cause Analysis",
         "Stack trace + failure log → classified failure type + impacted component."),
        ("🎛️", ACCENT4, "Regression Selection",
         "Changed file list → only the impacted test suites. No full-suite runs on every PR."),
        ("🏭", ACCENT2, "Provider Abstraction",
         "Swap LLM providers with one env var. CachingProvider + FallbackProvider built in."),
    ]

    x_start = Inches(0.6)
    y_start = Inches(1.72)
    card_w = Inches(4.0)
    card_h = Inches(1.42)
    gap    = Inches(0.2)

    for i, (icon, color, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = x_start + col * (card_w + gap)
        y = y_start + row * (card_h + gap)
        card(s, x, y, card_w, card_h, color, DARK_BG)
        add_textbox(s, icon,
                    x + Inches(0.15), y + Inches(0.1),
                    Inches(0.45), Inches(0.38), font_size=18)
        add_textbox(s, title,
                    x + Inches(0.6), y + Inches(0.1),
                    card_w - Inches(0.75), Inches(0.38),
                    font_size=13, bold=True, color=color)
        add_textbox(s, desc,
                    x + Inches(0.15), y + Inches(0.55),
                    card_w - Inches(0.3), Inches(0.78),
                    font_size=10, color=MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI CAPABILITIES (9 modules)
# ══════════════════════════════════════════════════════════════════════════════
def slide_09_ai_modules(prs):
    s = blank_slide(prs)
    eyebrow(s, "AI MODULES")
    slide_title(s, "Nine Independent Modules — One Shared Interface")
    add_divider(s, Inches(1.55))

    modules = [
        (ACCENT1, "TestCaseGenerator",
         "requirement → TestCase[]",
         "4–10 cases: positive, negative, edge, security"),
        (ACCENT2, "TestDataGenerator",
         "requirement → TestData",
         "Valid + invalid values, typed per page"),
        (ACCENT3, "ActionModelGenerator",
         "step description → ActionModel",
         "Structured JSON: goto / fill / click / assert"),
        (ACCENT4, "AssertionGenerator",
         "expectedResult + KB → assertion string",
         "Grounded in real selectors from knowledge base"),
        (ACCENT5, "SelfHealingLocator",
         "broken selector + page → healed selector",
         "Confidence score · KB-grounded · no hallucination"),
        (ACCENT1, "FlakyTestAnalyzer",
         "execution metrics → flakiness score",
         "Score 0–100 + specific fix recommendation"),
        (ACCENT3, "BugRootCauseAnalyzer",
         "error + trace + log → diagnosis",
         "Failure type · impacted component · confidence"),
        (ACCENT2, "CoverageAnalyzer",
         "requirements + tests → gap report",
         "Missing coverage · partial coverage · percentage"),
        (ACCENT4, "RegressionSelector",
         "changed files → impacted suites",
         "Filtered against catalog — no invented suite names"),
    ]

    cols = 3
    x_start = Inches(0.55)
    y_start = Inches(1.72)
    cw = Inches(4.12)
    ch = Inches(1.38)
    xg = Inches(0.12)
    yg = Inches(0.1)

    for i, (color, name, io, desc) in enumerate(modules):
        col = i % cols
        row = i // cols
        x = x_start + col * (cw + xg)
        y = y_start + row * (ch + yg)
        add_rect(s, x, y, Inches(0.05), ch, color)
        card(s, x + Inches(0.05), y, cw - Inches(0.05), ch, BORDER, DARK_BG)
        add_textbox(s, name,
                    x + Inches(0.2), y + Inches(0.1),
                    cw - Inches(0.3), Inches(0.34),
                    font_size=13, bold=True, color=color)
        add_textbox(s, io,
                    x + Inches(0.2), y + Inches(0.46),
                    cw - Inches(0.3), Inches(0.28),
                    font_size=10, color=WHITE, italic=True)
        add_textbox(s, desc,
                    x + Inches(0.2), y + Inches(0.78),
                    cw - Inches(0.3), Inches(0.48),
                    font_size=10, color=MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BEFORE vs AFTER
# ══════════════════════════════════════════════════════════════════════════════
def slide_10_before_after(prs):
    s = blank_slide(prs)
    eyebrow(s, "IMPACT")
    slide_title(s, "Before vs After")
    add_divider(s, Inches(1.55))

    # Before
    add_rect(s, Inches(0.6), Inches(1.72), Inches(5.8), Inches(5.2),
             RGBColor(0x20,0x0d,0x0d), RGBColor(0xda,0x36,0x33), Pt(1.5))
    add_textbox(s, "✗  BEFORE — Manual Process",
                Inches(0.8), Inches(1.8), Inches(5.4), Inches(0.38),
                font_size=11, bold=True, color=ACCENT5)
    before = [
        "Think through every scenario manually",
        "Write test cases in a document",
        "Create test data from scratch",
        "Translate steps to Playwright code",
        "Maintain selectors when UI changes",
        "Debug flaky tests reading stack traces",
        "No automated coverage visibility",
        "Full suite on every PR — 45 min wait",
    ]
    y = Inches(2.32)
    for item in before:
        add_textbox(s, f"→  {item}",
                    Inches(0.8), y, Inches(5.4), Inches(0.35),
                    font_size=11, color=MUTED)
        y += Inches(0.36)
    add_textbox(s, "⏱  Hours per feature · Breaks constantly",
                Inches(0.8), Inches(6.32), Inches(5.4), Inches(0.38),
                font_size=12, bold=True, color=ACCENT5)

    # After
    add_rect(s, Inches(6.8), Inches(1.72), Inches(5.8), Inches(5.2),
             RGBColor(0x0d,0x21,0x18), RGBColor(0x23,0x86,0x36), Pt(1.5))
    add_textbox(s, "✓  AFTER — AI-Powered Platform",
                Inches(7.0), Inches(1.8), Inches(5.4), Inches(0.38),
                font_size=11, bold=True, color=ACCENT2)
    after = [
        "One requirement → 8 test cases in 90 sec",
        "AI generates realistic test data",
        "Natural language → Playwright code",
        "Self-healing engine fixes broken selectors",
        "Flaky test analyzer diagnoses root causes",
        "AI explains every CI failure",
        "Coverage gap report on any commit",
        "CI runs only impacted tests",
    ]
    y = Inches(2.32)
    for item in after:
        add_textbox(s, f"→  {item}",
                    Inches(7.0), y, Inches(5.4), Inches(0.35),
                    font_size=11, color=WHITE)
        y += Inches(0.36)
    add_textbox(s, "⚡  90 seconds per feature · Self-maintaining",
                Inches(7.0), Inches(6.32), Inches(5.4), Inches(0.38),
                font_size=12, bold=True, color=ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
def slide_11_business_benefits(prs):
    s = blank_slide(prs)
    eyebrow(s, "BUSINESS VALUE")
    slide_title(s, "Why This Matters to the Organisation")
    add_divider(s, Inches(1.55))

    benefits = [
        (ACCENT1, "💰", "Faster Time-to-Market",
         "Test automation no longer blocks feature delivery. New spec files in minutes, not days."),
        (ACCENT2, "🛡️", "Higher Quality",
         "AI covers edge cases humans miss. Every requirement gets positive + negative + security tests."),
        (ACCENT3, "📈", "Scalable Coverage",
         "Add a new application in an hour. Knowledge base + config entry + one command."),
        (ACCENT4, "🔄", "Reduced Maintenance",
         "Self-healing eliminates most selector failures. Analysts spend time on strategy, not upkeep."),
        (ACCENT5, "👁️", "Full Visibility",
         "Coverage analyzer shows exactly which requirements lack tests. No more unknown blind spots."),
        (ACCENT2, "🔒", "Vendor Independence",
         "Swap LLM providers in one line. No lock-in to OpenAI, Google, or any single model."),
    ]

    x_start = Inches(0.6)
    y_start = Inches(1.72)
    bw = Inches(3.9)
    bh = Inches(1.65)
    gap = Inches(0.25)

    for i, (color, icon, title, desc) in enumerate(benefits):
        col = i % 3
        row = i // 3
        x = x_start + col * (bw + gap)
        y = y_start + row * (bh + gap)
        add_rect(s, x, y, bw, Inches(0.05), color)
        card(s, x, y + Inches(0.05), bw, bh - Inches(0.05), BORDER, DARK_BG)
        add_textbox(s, icon,
                    x + Inches(0.15), y + Inches(0.14),
                    Inches(0.45), Inches(0.42), font_size=20)
        add_textbox(s, title,
                    x + Inches(0.65), y + Inches(0.16),
                    bw - Inches(0.8), Inches(0.4),
                    font_size=14, bold=True, color=color)
        add_textbox(s, desc,
                    x + Inches(0.15), y + Inches(0.65),
                    bw - Inches(0.3), Inches(0.88),
                    font_size=11, color=MUTED, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TECHNICAL BENEFITS
# ══════════════════════════════════════════════════════════════════════════════
def slide_12_technical_benefits(prs):
    s = blank_slide(prs)
    eyebrow(s, "TECHNICAL EXCELLENCE")
    slide_title(s, "Enterprise-Grade Technical Foundation")
    add_divider(s, Inches(1.55))

    left_items = [
        (ACCENT1, "TypeScript Strict Mode",
         "No any types. No implicit undefined. Every module fully typed. 0 errors across the full codebase."),
        (ACCENT2, "104 Unit Tests",
         "The framework tests itself. Vitest runs all 104 tests in under one second using MockLLMProvider."),
        (ACCENT3, "ESM Modules",
         "Native Node.js ESM with nodenext resolution. Explicit .js extensions. Future-proof module system."),
        (ACCENT4, "Provider Abstraction",
         "Four-line LLMProvider interface. Swap Gemini for GPT-5 next year — zero code changes in modules."),
    ]
    right_items = [
        (ACCENT2, "CachingLLMProvider",
         "File-based response cache. Development runs hit disk, not the API. No repeated API cost."),
        (ACCENT5, "FallbackProvider + Circuit Breaker",
         "Auto-switch on failure. Trips after 5 consecutive errors. Prevents retry storms."),
        (ACCENT1, "AIJsonParser",
         "Strips markdown fences from LLM responses. Parses as typed JSON. Every module uses this."),
        (ACCENT3, "Allure + Playwright Reporting",
         "Dual reporting: Playwright HTML report + Allure with categories, trends, and timeline."),
    ]

    y = Inches(1.75)
    for (lc, lt, ld), (rc, rt, rd) in zip(left_items, right_items):
        add_rect(s, Inches(0.6), y + Inches(0.1),
                 Inches(0.04), Inches(0.6), lc)
        add_textbox(s, lt, Inches(0.78), y,
                    Inches(5.5), Inches(0.36),
                    font_size=13, bold=True, color=lc)
        add_textbox(s, ld, Inches(0.78), y + Inches(0.36),
                    Inches(5.5), Inches(0.5),
                    font_size=11, color=MUTED, wrap=True)

        add_rect(s, Inches(6.9), y + Inches(0.1),
                 Inches(0.04), Inches(0.6), rc)
        add_textbox(s, rt, Inches(7.08), y,
                    Inches(5.5), Inches(0.36),
                    font_size=13, bold=True, color=rc)
        add_textbox(s, rd, Inches(7.08), y + Inches(0.36),
                    Inches(5.5), Inches(0.5),
                    font_size=11, color=MUTED, wrap=True)
        y += Inches(1.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — REPORTS GENERATED
# ══════════════════════════════════════════════════════════════════════════════
def slide_13_reports(prs):
    s = blank_slide(prs)
    eyebrow(s, "REPORTING")
    slide_title(s, "Reports Generated by the Framework")
    add_divider(s, Inches(1.55))

    reports = [
        (ACCENT1, "Playwright HTML Report",
         "npm run report",
         "Pass/fail counts · test durations · screenshot on failure · retry history · test steps",
         "reports/latest/playwright/"),
        (ACCENT4, "Allure Report",
         "npm run allure:full",
         "Test categories · timeline · trends · severity · owner · defects dashboard",
         "reports/latest/allure/report/"),
        (ACCENT2, "Coverage Gap Report",
         "npm run demo:coverage",
         "Uncovered requirements · partial coverage · coverage % · gap recommendations",
         "Printed to console + JSON output"),
        (ACCENT3, "Flaky Test Report",
         "npm run demo:flaky",
         "Flakiness score 0–100 · root cause type · specific fix recommendation per test",
         "Printed to console"),
        (ACCENT5, "Root Cause Report",
         "npm run demo:rootcause",
         "Failure type · impacted component · confidence · recommended fix action",
         "Printed to console"),
        (ACCENT1, "Regression Impact Report",
         "npm run demo:regression",
         "Which suites are impacted by changed files · why each suite was selected",
         "Printed to console"),
    ]

    x_start = Inches(0.6)
    y_start = Inches(1.72)
    rw = Inches(3.9)
    rh = Inches(1.55)
    gap = Inches(0.24)

    for i, (color, title, cmd, desc, path) in enumerate(reports):
        col = i % 3
        row = i // 3
        x = x_start + col * (rw + gap)
        y = y_start + row * (rh + gap)
        card(s, x, y, rw, rh, color, DARK_BG)
        add_textbox(s, title, x + Inches(0.15), y + Inches(0.1),
                    rw - Inches(0.3), Inches(0.36),
                    font_size=12, bold=True, color=color)
        add_textbox(s, f"$ {cmd}", x + Inches(0.15), y + Inches(0.48),
                    rw - Inches(0.3), Inches(0.28),
                    font_size=10, color=ACCENT2, italic=True)
        add_textbox(s, desc, x + Inches(0.15), y + Inches(0.78),
                    rw - Inches(0.3), Inches(0.5),
                    font_size=10, color=MUTED, wrap=True)
        add_textbox(s, path, x + Inches(0.15), y + rh - Inches(0.3),
                    rw - Inches(0.3), Inches(0.25),
                    font_size=9, color=BORDER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — REUSABILITY
# ══════════════════════════════════════════════════════════════════════════════
def slide_14_reusability(prs):
    s = blank_slide(prs)
    eyebrow(s, "REUSABILITY")
    slide_title(s, "Adding a New Application — 4 Steps")
    add_divider(s, Inches(1.55))

    add_textbox(s, "The entire framework is application-agnostic. Any web application can be onboarded in under an hour.",
                Inches(0.6), Inches(1.72), Inches(12.1), Inches(0.45),
                font_size=13, color=MUTED)

    steps = [
        (ACCENT1, "1", "Add Knowledge Base JSON",
         "Create pipeline/kb/pages/my-page.json with pageName, url, selectors, messages.",
         "Or run: npm run kb:generate -- --url https://myapp.com/login",
         "Time: 5 minutes (or 30 seconds with KnowledgeBaseGenerator)"),
        (ACCENT3, "2", "Add Suite to platform.json",
         'Add one entry: { "name": "My Page", "page": "my-page", "outputFile": "my-page.spec.ts" }',
         "File: config/platform.json",
         "Time: 2 minutes"),
        (ACCENT2, "3", "Generate Tests",
         "npm run generate:all — creates tests/pages/MyPage.ts and tests/e2e/my-page.spec.ts",
         "Both files follow the same conventions as existing hand-written tests.",
         "Time: 90 seconds"),
        (ACCENT4, "4", "Enrich the Page Object",
         "Add behavior methods (login(), search(), addToCart()) to the generated POM.",
         "Follow the pattern in tests/pages/AeHomePage.ts",
         "Time: 30 minutes for a typical page"),
    ]

    y = Inches(2.3)
    for color, num, title, desc, detail, timing in steps:
        add_rect(s, Inches(0.6), y, Inches(0.5), Inches(1.12), color)
        add_textbox(s, num,
                    Inches(0.6), y + Inches(0.35), Inches(0.5), Inches(0.44),
                    font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        card(s, Inches(1.15), y, Inches(11.5), Inches(1.12), color, DARK_BG)
        add_textbox(s, title,
                    Inches(1.35), y + Inches(0.08), Inches(11), Inches(0.36),
                    font_size=14, bold=True, color=color)
        add_textbox(s, desc,
                    Inches(1.35), y + Inches(0.46), Inches(9.5), Inches(0.3),
                    font_size=11, color=WHITE)
        add_textbox(s, detail,
                    Inches(1.35), y + Inches(0.76), Inches(9.0), Inches(0.28),
                    font_size=10, color=MUTED, italic=True)
        add_textbox(s, timing,
                    Inches(10.4), y + Inches(0.76), Inches(2.15), Inches(0.28),
                    font_size=10, color=ACCENT2, align=PP_ALIGN.RIGHT)
        y += Inches(1.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ENTERPRISE READINESS
# ══════════════════════════════════════════════════════════════════════════════
def slide_15_enterprise(prs):
    s = blank_slide(prs)
    eyebrow(s, "ENTERPRISE READINESS")
    slide_title(s, "Production-Ready From Day One")
    add_divider(s, Inches(1.55))

    columns = [
        ("CI/CD Integration", ACCENT1, [
            "4-job GitHub Actions pipeline",
            "TypeScript → ESLint → Unit → Smoke",
            "Smoke tests gate on type check + unit tests",
            "Artifacts uploaded with 7-day retention",
            "PR checks enforce 0 TS errors",
        ]),
        ("Multi-Environment Support", ACCENT2, [
            "config/environments/qa.env",
            "config/environments/uat.env",
            "config/environments/production.env",
            "npm run test:qa / test:uat / test:prod",
            "BASE_URL injected per environment",
        ]),
        ("Reliability Architecture", ACCENT3, [
            "FallbackProvider with circuit breaker",
            "Trips after 5 consecutive LLM failures",
            "CachingProvider skips repeated API calls",
            "MockLLMProvider for offline development",
            "Retry + backoff on transient failures",
        ]),
        ("Code Quality Standards", ACCENT4, [
            "TypeScript strict — no any, no unknown",
            "ESLint + Playwright rules enforced",
            "Prettier formatting across all files",
            "104 vitest unit tests on framework code",
            "Every PR: typecheck + lint + unit tests",
        ]),
    ]

    x = Inches(0.6)
    cw = Inches(2.95)
    for title, color, items in columns:
        add_rect(s, x, Inches(1.72), cw, Inches(0.05), color)
        card(s, x, Inches(1.77), cw, Inches(5.2), BORDER, DARK_BG)
        add_textbox(s, title,
                    x + Inches(0.15), Inches(1.88),
                    cw - Inches(0.3), Inches(0.4),
                    font_size=13, bold=True, color=color)
        y = Inches(2.38)
        for item in items:
            add_textbox(s, f"◆  {item}",
                        x + Inches(0.15), y,
                        cw - Inches(0.3), Inches(0.38),
                        font_size=10, color=WHITE, wrap=True)
            y += Inches(0.38)
        x += cw + Inches(0.18)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — FUTURE ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
def slide_16_roadmap(prs):
    s = blank_slide(prs)
    eyebrow(s, "ROADMAP")
    slide_title(s, "What's Next — Planned Enhancements")
    add_divider(s, Inches(1.55))

    add_textbox(s, "The following are planned features, not yet implemented.",
                Inches(0.6), Inches(1.72), Inches(12), Inches(0.38),
                font_size=12, color=ACCENT3, italic=True)

    phases = [
        ("Phase 1 — Q3 2026", ACCENT2, [
            "GherkinGenerator — export test cases as Gherkin .feature files",
            "Visual regression integration (screenshot diffing on generated tests)",
            "Slack/Teams notification on coverage gap findings",
        ]),
        ("Phase 2 — Q4 2026", ACCENT1, [
            "Multi-page orchestration — cross-page user journey generation",
            "API contract test generation alongside UI tests",
            "Self-updating knowledge base (auto-detect selector drift)",
        ]),
        ("Phase 3 — 2027", ACCENT4, [
            "Agent loop for test exploration (autonomous scenario discovery)",
            "Integration with Jira/Linear for requirement sync",
            "LLM fine-tuning on internal codebase patterns",
        ]),
    ]

    y = Inches(2.28)
    for phase, color, items in phases:
        add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.38), color)
        add_textbox(s, phase,
                    Inches(0.8), y + Inches(0.05),
                    Inches(4), Inches(0.3),
                    font_size=12, bold=True, color=WHITE)
        y += Inches(0.38)
        for item in items:
            card(s, Inches(0.6), y, Inches(12.1), Inches(0.5), BORDER, DARK_BG)
            add_textbox(s, f"○  {item}",
                        Inches(0.85), y + Inches(0.1),
                        Inches(11.5), Inches(0.32),
                        font_size=12, color=WHITE)
            y += Inches(0.56)
        y += Inches(0.1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — LIVE DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
def slide_17_demo_flow(prs):
    s = blank_slide(prs)
    eyebrow(s, "LIVE DEMO")
    slide_title(s, "Demo Script — What to Run and When")
    add_divider(s, Inches(1.55))

    demos = [
        (ACCENT2, "1. Unit Tests",         "npm run test:unit",
         "Show 104 tests pass in < 1 sec. Framework tests itself."),
        (ACCENT1, "2. TypeScript Check",   "npm run typecheck",
         "0 errors. Strict mode. Every module typed."),
        (ACCENT3, "3. Full Pipeline Demo", "npm run demo",
         "Requirement → test cases → data → action model → .spec.ts"),
        (ACCENT5, "4. Self-Healing",       "npm run demo:healing",
         "Broken selector → healed selector + confidence score."),
        (ACCENT4, "5. Flaky Analyzer",     "npm run demo:flaky",
         "Execution metrics → flakiness score + fix recommendation."),
        (ACCENT2, "6. Root Cause",         "npm run demo:rootcause",
         "Stack trace + log → classified failure + component."),
        (ACCENT1, "7. Coverage Gaps",      "npm run demo:coverage",
         "Requirements vs tests → missing coverage report."),
        (ACCENT3, "8. Run Smoke Tests",    "npm run test:smoke",
         "Live @smoke suite against automationexercise.com."),
    ]

    y = Inches(1.72)
    for i, (color, title, cmd, desc) in enumerate(demos):
        add_rect(s, Inches(0.6), y + Inches(0.08),
                 Inches(0.04), Inches(0.52), color)
        add_textbox(s, title,
                    Inches(0.78), y, Inches(2.8), Inches(0.36),
                    font_size=12, bold=True, color=color)
        add_textbox(s, f"$ {cmd}",
                    Inches(3.7), y, Inches(4.2), Inches(0.36),
                    font_size=11, color=ACCENT2)
        add_textbox(s, desc,
                    Inches(8.1), y, Inches(5.0), Inches(0.36),
                    font_size=11, color=MUTED)
        y += Inches(0.66)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Q&A
# ══════════════════════════════════════════════════════════════════════════════
def slide_18_qa(prs):
    s = blank_slide(prs)
    eyebrow(s, "Q & A")
    slide_title(s, "Common Questions")
    add_divider(s, Inches(1.55))

    qas = [
        (ACCENT1,
         "Q: How do we know the generated tests are correct?",
         "A: AI modules validate their own output (min 4 test cases enforced). "
            "Generated tests go through the same PR review as hand-written ones. "
            "The AI describes intent — a rule-based renderer writes the actual code."),
        (ACCENT3,
         "Q: What happens when the AI returns a wrong answer?",
         "A: AIJsonParser validates structure before any code runs. Knowledge base "
            "grounds selectors in reality — the AI cannot invent ones that don't exist. "
            "MockLLMProvider lets us test all module logic without any AI dependency."),
        (ACCENT4,
         "Q: Why not use an agent loop instead of a fixed pipeline?",
         "A: Predictability over flexibility. Test generation benefits from consistent "
            "output structure. We documented the full tradeoff in docs/ai-architecture-evaluation.md."),
        (ACCENT2,
         "Q: What does it take to add our application?",
         "A: KB JSON + one config entry + npm run generate:all. "
            "Under an hour for a typical page. See docs/ONBOARDING.md for the step-by-step guide."),
    ]

    y = Inches(1.72)
    for color, q, a in qas:
        card(s, Inches(0.6), y, Inches(12.1), Inches(1.2), color, DARK_BG)
        add_textbox(s, q,
                    Inches(0.8), y + Inches(0.08),
                    Inches(11.7), Inches(0.36),
                    font_size=12, bold=True, color=color)
        add_textbox(s, a,
                    Inches(0.8), y + Inches(0.46),
                    Inches(11.7), Inches(0.64),
                    font_size=11, color=MUTED, wrap=True)
        y += Inches(1.28)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
def slide_19_thankyou(prs):
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
             RGBColor(0x0a, 0x14, 0x20))

    # Accent bar
    add_rect(s, Inches(0), Inches(2.5), SLIDE_W, Inches(0.06), ACCENT1)

    add_textbox(s, "THANK YOU",
                Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
                font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
    add_textbox(s,
        "AI-Powered Playwright\nTest Generation Framework",
                Inches(1.5), Inches(1.65), Inches(10), Inches(1.5),
                font_size=38, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, wrap=True)

    add_textbox(s,
        "Scan the QR code, read docs/DEMO_SCRIPT.md, or pair with me to onboard your project.",
                Inches(1.5), Inches(3.3), Inches(10), Inches(0.55),
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

    links = [
        ("Docs",          "docs/GETTING-STARTED.md",   ACCENT1),
        ("Onboarding",    "docs/ONBOARDING.md",         ACCENT2),
        ("Commands",      "docs/COMMANDS.md",           ACCENT3),
        ("Architecture",  "docs/FRAMEWORK_STRUCTURE.md",ACCENT4),
    ]
    x = Inches(2.0)
    lw = Inches(2.15)
    for label, ref, color in links:
        card(s, x, Inches(4.2), lw, Inches(0.8), color, DARK_BG)
        add_textbox(s, label, x + Inches(0.1), Inches(4.26),
                    lw - Inches(0.2), Inches(0.3),
                    font_size=11, bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, ref, x + Inches(0.05), Inches(4.58),
                    lw - Inches(0.1), Inches(0.28),
                    font_size=9, color=MUTED,
                    align=PP_ALIGN.CENTER)
        x += lw + Inches(0.25)

    add_textbox(s, "Ravi Gaygol  ·  rgaygol@rxsense.com  ·  2026",
                Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
                font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs);           print("  ✓ 01 Title")
    slide_02_exec_summary(prs);    print("  ✓ 02 Executive Summary")
    slide_03_problem(prs);         print("  ✓ 03 Business Problem")
    slide_04_solution(prs);        print("  ✓ 04 Solution Overview")
    slide_05_architecture(prs);    print("  ✓ 05 Architecture")
    slide_06_workflow(prs);        print("  ✓ 06 AI Workflow")
    slide_07_folder_structure(prs);print("  ✓ 07 Folder Structure")
    slide_08_features(prs);        print("  ✓ 08 Key Features")
    slide_09_ai_modules(prs);      print("  ✓ 09 AI Modules")
    slide_10_before_after(prs);    print("  ✓ 10 Before vs After")
    slide_11_business_benefits(prs);print(" ✓ 11 Business Benefits")
    slide_12_technical_benefits(prs);print(" ✓ 12 Technical Benefits")
    slide_13_reports(prs);         print("  ✓ 13 Reports")
    slide_14_reusability(prs);     print("  ✓ 14 Reusability")
    slide_15_enterprise(prs);      print("  ✓ 15 Enterprise Readiness")
    slide_16_roadmap(prs);         print("  ✓ 16 Roadmap")
    slide_17_demo_flow(prs);       print("  ✓ 17 Demo Flow")
    slide_18_qa(prs);              print("  ✓ 18 Q&A")
    slide_19_thankyou(prs);        print("  ✓ 19 Thank You")

    out = "docs/AI-Platform-Presentation.pptx"
    prs.save(out)
    print(f"\n✅  Saved → {out}")
    print(f"   19 slides · {SLIDE_W.inches:.0f}\" × {SLIDE_H.inches:.1f}\"")

if __name__ == "__main__":
    main()
